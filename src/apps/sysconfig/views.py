import io
import os

from django.conf import settings
from django.contrib.auth.decorators import login_required
from django.http import HttpResponse, Http404
from django.shortcuts import render
from core.permissions.roles import IsAdmin


@login_required
def system_view(request):
    """시스템 설정 페이지 (관리자 전용)"""
    from apps.schools.models import SupportCenter
    centers = list(SupportCenter.objects.filter(is_active=True).order_by('id').values('id', 'name'))
    return render(request, 'sysconfig/index.html', {'centers': centers})


@login_required
def guide_view(request):
    """사용 안내 페이지 — MODULE_REGISTRY 기반 자동 생성"""
    from core.guide import get_guide_modules
    modules = get_guide_modules()
    return render(request, 'guide/index.html', {'modules': modules})


@login_required
def guide_pptx_view(request, module_key):
    """사용 안내 모듈별 PPTX 생성 (슈퍼어드민 전용)"""
    if request.user.role != 'superadmin':
        raise Http404

    from core.guide import get_guide_modules
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    # 모듈 데이터 조회
    modules = {m['key']: m for m in get_guide_modules()}
    mod = modules.get(module_key)
    if not mod:
        raise Http404

    # ── 슬라이드 설정 (Widescreen 33.87 × 19.05 cm) ──────────────────────────
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(slide_layout)

    W = prs.slide_width
    H = prs.slide_height

    # ── 폰트 경로 ────────────────────────────────────────────────────────────
    FONT_DIR   = os.path.join(settings.BASE_DIR, 'static', 'fonts')
    FONT_REG   = os.path.join(FONT_DIR, 'malgun.ttf')
    FONT_BOLD  = os.path.join(FONT_DIR, 'malgunbd.ttf')

    def _set_font(run, size_pt, bold=False, color=None):
        run.font.name = '맑은 고딕'
        run.font.size = Pt(size_pt)
        run.font.bold = bold
        if color:
            run.font.color.rgb = RGBColor(*color)

    def _add_tf(slide, left, top, width, height, word_wrap=True):
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = word_wrap
        return tf

    # ── 색상 팔레트 ───────────────────────────────────────────────────────────
    COLOR_TITLE_BG  = RGBColor(0x1E, 0x40, 0xAF)   # 진파랑
    COLOR_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
    COLOR_DARK      = RGBColor(0x1E, 0x29, 0x3B)
    COLOR_GRAY      = RGBColor(0x64, 0x74, 0x8B)
    COLOR_ACCENT    = RGBColor(0x0E, 0xA5, 0xE9)
    COLOR_SEC_BG    = RGBColor(0xF1, 0xF5, 0xF9)
    COLOR_BORDER    = RGBColor(0xCB, 0xD5, 0xE1)

    MARGIN   = Inches(0.25)
    TITLE_H  = Inches(0.75)

    # ── 제목 배너 ─────────────────────────────────────────────────────────────
    title_bg = slide.shapes.add_shape(
        1, 0, 0, W, TITLE_H  # MSO_SHAPE_TYPE.RECTANGLE
    )
    title_bg.fill.solid()
    title_bg.fill.fore_color.rgb = COLOR_TITLE_BG
    title_bg.line.fill.background()

    tf_title = _add_tf(slide, MARGIN, Emu(0), W - MARGIN * 2, TITLE_H)
    p = tf_title.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = mod['label']
    _set_font(run, 22, bold=True, color=(0xFF, 0xFF, 0xFF))

    # ── 레이아웃 치수 ─────────────────────────────────────────────────────────
    BODY_TOP   = TITLE_H + MARGIN
    BODY_H     = H - BODY_TOP - MARGIN
    LEFT_W     = Inches(5.8)
    RIGHT_W    = W - LEFT_W - MARGIN * 3
    LEFT_LEFT  = MARGIN
    RIGHT_LEFT = LEFT_LEFT + LEFT_W + MARGIN

    # ── 좌측 패널 배경 ────────────────────────────────────────────────────────
    left_bg = slide.shapes.add_shape(1, LEFT_LEFT, BODY_TOP, LEFT_W, BODY_H)
    left_bg.fill.solid()
    left_bg.fill.fore_color.rgb = COLOR_SEC_BG
    left_bg.line.color.rgb = COLOR_BORDER

    # ── 좌측: 이미지 ──────────────────────────────────────────────────────────
    IMG_H      = Inches(3.5)
    img_path   = os.path.join(settings.BASE_DIR, 'static', 'img', 'guide', f'{module_key}.png')
    img_left   = LEFT_LEFT + MARGIN
    img_top    = BODY_TOP + MARGIN
    img_width  = LEFT_W - MARGIN * 2

    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, img_left, img_top, img_width, IMG_H)
    else:
        # placeholder
        ph = slide.shapes.add_shape(1, img_left, img_top, img_width, IMG_H)
        ph.fill.solid()
        ph.fill.fore_color.rgb = RGBColor(0xE2, 0xE8, 0xF0)
        ph.line.color.rgb = COLOR_BORDER
        tf_ph = ph.text_frame
        tf_ph.text = '스크린샷 준비 중'
        run = tf_ph.paragraphs[0].runs[0]
        _set_font(run, 12, color=(0x94, 0xA3, 0xB8))
        tf_ph.paragraphs[0].alignment = PP_ALIGN.CENTER

    # ── 좌측: 모듈 설명 ───────────────────────────────────────────────────────
    desc_top   = img_top + IMG_H + MARGIN * 0.5
    desc_h     = BODY_TOP + BODY_H - desc_top - MARGIN
    tf_desc    = _add_tf(slide, img_left, desc_top, img_width, desc_h)
    p = tf_desc.paragraphs[0]
    run = p.add_run()
    run.text = mod['summary']
    _set_font(run, 9, color=(0x1E, 0x29, 0x3B))

    # ── 우측 패널 배경 ────────────────────────────────────────────────────────
    right_bg = slide.shapes.add_shape(1, RIGHT_LEFT, BODY_TOP, RIGHT_W, BODY_H)
    right_bg.fill.solid()
    right_bg.fill.fore_color.rgb = COLOR_WHITE
    right_bg.line.color.rgb = COLOR_BORDER

    cur_top = BODY_TOP + MARGIN * 0.8

    # ── 우측: 탭 구성 ─────────────────────────────────────────────────────────
    if mod['tabs']:
        tf_sec = _add_tf(slide, RIGHT_LEFT + MARGIN * 0.5, cur_top, RIGHT_W - MARGIN, Inches(0.3))
        p = tf_sec.paragraphs[0]
        run = p.add_run()
        run.text = '▶ 탭 구성'
        _set_font(run, 10, bold=True, color=(0x0E, 0xA5, 0xE9))
        cur_top += Inches(0.32)

        for tab in mod['tabs']:
            tab_h = Inches(0.55)
            tab_bg = slide.shapes.add_shape(
                1,
                RIGHT_LEFT + MARGIN * 0.5,
                cur_top,
                RIGHT_W - MARGIN,
                tab_h,
            )
            tab_bg.fill.solid()
            tab_bg.fill.fore_color.rgb = RGBColor(0xEF, 0xF6, 0xFF)
            tab_bg.line.color.rgb = RGBColor(0xBF, 0xDB, 0xFE)

            tf_tab = _add_tf(
                slide,
                RIGHT_LEFT + MARGIN * 0.7,
                cur_top + Emu(30000),
                RIGHT_W - MARGIN * 1.4,
                tab_h,
            )
            p0 = tf_tab.paragraphs[0]
            r0 = p0.add_run()
            r0.text = tab['name']
            _set_font(r0, 8.5, bold=True, color=(0x1E, 0x40, 0xAF))

            p1 = tf_tab.add_paragraph()
            r1 = p1.add_run()
            r1.text = tab['desc'][:80] + ('…' if len(tab['desc']) > 80 else '')
            _set_font(r1, 7.5, color=(0x47, 0x55, 0x69))

            cur_top += tab_h + Emu(40000)

    cur_top += Inches(0.1)

    # ── 우측: 주요 기능 ───────────────────────────────────────────────────────
    if mod['features']:
        tf_sec2 = _add_tf(slide, RIGHT_LEFT + MARGIN * 0.5, cur_top, RIGHT_W - MARGIN, Inches(0.3))
        p = tf_sec2.paragraphs[0]
        run = p.add_run()
        run.text = '▶ 주요 기능'
        _set_font(run, 10, bold=True, color=(0x0E, 0xA5, 0xE9))
        cur_top += Inches(0.32)

        feat_h = BODY_TOP + BODY_H - cur_top - MARGIN
        tf_feat = _add_tf(slide, RIGHT_LEFT + MARGIN * 0.5, cur_top, RIGHT_W - MARGIN, feat_h)
        for i, feat in enumerate(mod['features']):
            p = tf_feat.paragraphs[0] if i == 0 else tf_feat.add_paragraph()
            run = p.add_run()
            run.text = f'✓  {feat}'
            _set_font(run, 8.5, color=(0x1E, 0x29, 0x3B))
            p.space_after = Pt(3)

    # ── 파일 응답 ─────────────────────────────────────────────────────────────
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)

    from urllib.parse import quote
    filename = f'{mod["label"]}.pptx'
    encoded  = quote(filename, safe='')
    response = HttpResponse(
        buf,
        content_type='application/vnd.openxmlformats-officedocument.presentationml.presentation',
    )
    response['Content-Disposition'] = f"attachment; filename*=UTF-8''{encoded}"
    return response

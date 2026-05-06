from django.shortcuts import render
from django.contrib.auth.decorators import login_required
from rest_framework import viewsets, status
from rest_framework.decorators import action
from rest_framework.response import Response
from rest_framework.permissions import IsAuthenticated
from django.utils import timezone
from django.db.models import Count, Q

@login_required
def network_monitor_view(request):
    return render(request, 'network/monitor.html')


@login_required
def ap_analyzer_view(request):
    return render(request, 'ap_analyzer/index.html')


from .models import (
    NetworkDevice, NetworkPort, NetworkLink,
    NetworkTopology, NetworkEvent, SnmpDevice, SnmpMetric, NetworkCommand
)
from .serializers import (
    NetworkDeviceListSerializer, NetworkDeviceDetailSerializer,
    NetworkPortSerializer, NetworkTopologySerializer,
    NetworkEventSerializer, SnmpMetricSerializer, NetworkCommandSerializer
)
from core.permissions.roles import IsAdmin
from core.pagination import StandardPagination


class NetworkDeviceViewSet(viewsets.ModelViewSet):
    """네트워크 장비 관리 (NMS)"""
    permission_classes = [IsAuthenticated]
    pagination_class = StandardPagination

    def get_serializer_class(self):
        if self.action == 'retrieve':
            return NetworkDeviceDetailSerializer
        return NetworkDeviceListSerializer

    def get_queryset(self):
        qs = NetworkDevice.objects.select_related('school', 'asset').order_by('school', 'name')
        school_id = self.request.query_params.get('school_id')
        if school_id:
            qs = qs.filter(school_id=school_id)
        device_type = self.request.query_params.get('device_type')
        if device_type:
            qs = qs.filter(device_type=device_type)
        st = self.request.query_params.get('status')
        if st:
            qs = qs.filter(status=st)
        q = self.request.query_params.get('q')
        if q:
            qs = qs.filter(
                Q(name__icontains=q) | Q(ip_address__icontains=q) |
                Q(hostname__icontains=q) | Q(serial_number__icontains=q)
            )
        return qs

    @action(detail=False, methods=['get'])
    def summary(self, request):
        """장비 현황 요약"""
        qs = NetworkDevice.objects.values('status').annotate(cnt=Count('id'))
        status_counts = {item['status']: item['cnt'] for item in qs}
        down_devices = NetworkDevice.objects.filter(status='down').values(
            'id', 'name', 'ip_address', 'school__name'
        )[:10]
        return Response({
            'total':    NetworkDevice.objects.count(),
            'up':       status_counts.get('up', 0),
            'down':     status_counts.get('down', 0),
            'warning':  status_counts.get('warning', 0),
            'unknown':  status_counts.get('unknown', 0),
            'down_devices': list(down_devices),
        })

    @action(detail=True, methods=['get'])
    def ports(self, request, pk=None):
        """장비 포트 목록"""
        device = self.get_object()
        ports = device.ports.all()
        return Response(NetworkPortSerializer(ports, many=True).data)

    @action(detail=True, methods=['post'])
    def execute_command(self, request, pk=None):
        """원격 명령 실행"""
        device = self.get_object()
        command_type = request.data.get('command_type')
        command      = request.data.get('command', '')
        if not command_type:
            return Response({'error': '명령 유형이 필요합니다.'}, status=status.HTTP_400_BAD_REQUEST)
        cmd = NetworkCommand.objects.create(
            device=device,
            command_type=command_type,
            command=command,
            executed_by=request.user,
            status='pending',
        )
        # 비동기 실행은 tasks.py에서 처리 (Celery)
        from .tasks import execute_network_command
        execute_network_command.delay(cmd.id)
        return Response(NetworkCommandSerializer(cmd).data, status=status.HTTP_201_CREATED)

    @action(detail=True, methods=['get'])
    def metrics(self, request, pk=None):
        """SNMP 수집 지표"""
        device = self.get_object()
        metric_name = request.query_params.get('metric')
        qs = device.snmp_metrics.order_by('-collected_at')
        if metric_name:
            qs = qs.filter(metric_name=metric_name)
        qs = qs[:100]
        return Response(SnmpMetricSerializer(qs, many=True).data)

    @action(detail=False, methods=['get'])
    def school_status(self, request):
        """학교별 장비 상태 현황"""
        school_id = request.query_params.get('school_id')
        if not school_id:
            return Response({'error': 'school_id가 필요합니다.'}, status=status.HTTP_400_BAD_REQUEST)
        devices = NetworkDevice.objects.filter(school_id=school_id)
        total   = devices.count()
        up      = devices.filter(status='up').count()
        down    = devices.filter(status='down').count()
        return Response({
            'school_id': school_id,
            'total': total, 'up': up, 'down': down,
            'warning': devices.filter(status='warning').count(),
            'devices': NetworkDeviceListSerializer(devices, many=True).data,
        })


class NetworkTopologyViewSet(viewsets.ReadOnlyModelViewSet):
    """네트워크 토폴로지 스냅샷"""
    serializer_class = NetworkTopologySerializer
    permission_classes = [IsAuthenticated]
    pagination_class = StandardPagination

    def get_queryset(self):
        qs = NetworkTopology.objects.select_related('school')
        school_id = self.request.query_params.get('school_id')
        if school_id:
            qs = qs.filter(school_id=school_id)
        return qs

    @action(detail=False, methods=['get'])
    def latest(self, request):
        """학교별 최신 토폴로지"""
        school_id = request.query_params.get('school_id')
        if not school_id:
            return Response({'error': 'school_id가 필요합니다.'}, status=status.HTTP_400_BAD_REQUEST)
        topo = NetworkTopology.objects.filter(school_id=school_id).order_by('-scanned_at').first()
        if not topo:
            return Response({'error': '토폴로지 데이터가 없습니다.'}, status=status.HTTP_404_NOT_FOUND)
        return Response(NetworkTopologySerializer(topo).data)

    @action(detail=False, methods=['post'])
    def generate(self, request):
        """현재 장비/링크 데이터로 토폴로지 스냅샷 생성"""
        school_id = request.data.get('school_id')
        if not school_id:
            return Response({'error': 'school_id가 필요합니다.'}, status=status.HTTP_400_BAD_REQUEST)
        try:
            from .services import generate_and_save_topology
            topo = generate_and_save_topology(school_id)
            return Response(NetworkTopologySerializer(topo).data, status=status.HTTP_201_CREATED)
        except Exception as exc:
            return Response({'error': str(exc)}, status=status.HTTP_400_BAD_REQUEST)

    @action(detail=False, methods=['get'])
    def live(self, request):
        """저장 없이 현재 장비/링크 데이터를 실시간으로 반환"""
        school_id = request.query_params.get('school_id')
        if not school_id:
            return Response({'error': 'school_id가 필요합니다.'}, status=status.HTTP_400_BAD_REQUEST)
        from .services import build_topology_data
        return Response(build_topology_data(school_id))

    @action(detail=False, methods=['post'])
    def import_json(self, request):
        """Claude 분석 JSON → NetworkDevice + NetworkLink 저장"""
        from apps.schools.models import School
        from .models import NetworkDevice, NetworkLink

        school_id = request.data.get('school_id')
        data      = request.data.get('data', {})
        if not school_id or not data:
            return Response({'error': 'school_id와 data가 필요합니다.'}, status=status.HTTP_400_BAD_REQUEST)

        try:
            school = School.objects.get(id=school_id)
        except School.DoesNotExist:
            return Response({'error': '학교를 찾을 수 없습니다.'}, status=status.HTTP_404_NOT_FOUND)

        nodes = data.get('nodes', [])
        edges = data.get('edges', [])

        # 기존 장비/링크 초기화 (재등록)
        NetworkLink.objects.filter(from_device__school=school, link_type='manual').delete()
        NetworkDevice.objects.filter(school=school, ip_address__isnull=True, snmp_enabled=False).delete()

        name_to_device = {}
        created_devices = 0
        for node in nodes:
            name = node.get('name', '').strip()
            if not name:
                continue
            dev, created = NetworkDevice.objects.get_or_create(
                school=school, name=name,
                defaults={
                    'device_type':  node.get('device_type', 'switch'),
                    'model':        node.get('model', ''),
                    'location':     node.get('location', ''),
                    'network_type': node.get('network_type', ''),
                    'status':       'unknown',
                },
            )
            if created:
                created_devices += 1
            name_to_device[name] = dev

        created_links = 0
        CABLE_MAP = {'광': 'fiber', 'Cat6': 'cat6', 'Cat5e': 'cat5e', 'Cat5': 'cat5'}
        for edge in edges:
            fd = name_to_device.get(edge.get('from', ''))
            td = name_to_device.get(edge.get('to', ''))
            if fd and td and fd != td:
                cable_raw = edge.get('cable_type', '')
                NetworkLink.objects.create(
                    from_device=fd, to_device=td,
                    link_type='manual', is_active=True,
                    cable_type=CABLE_MAP.get(cable_raw, 'unknown'),
                    network_type=edge.get('network_type', ''),
                )
                created_links += 1

        return Response({
            'school': school.name,
            'created_devices': created_devices,
            'created_links':   created_links,
        })

    @action(detail=False, methods=['get'])
    def export_csv(self, request):
        """학교 장비 목록 CSV 다운로드"""
        import csv
        from django.http import HttpResponse
        from apps.schools.models import School
        from .models import NetworkDevice

        school_id = request.query_params.get('school_id')
        if not school_id:
            return Response({'error': 'school_id가 필요합니다.'}, status=status.HTTP_400_BAD_REQUEST)

        try:
            school = School.objects.get(id=school_id)
        except School.DoesNotExist:
            return Response({'error': '학교를 찾을 수 없습니다.'}, status=status.HTTP_404_NOT_FOUND)

        devices = NetworkDevice.objects.filter(school=school).order_by('network_type', 'name')

        from datetime import datetime
        import urllib.parse
        today = datetime.now().strftime('%Y%m%d')
        response = HttpResponse(content_type='text/csv; charset=utf-8-sig')
        filename = f'토폴로지_{school.name}_{today}.csv'
        encoded = urllib.parse.quote(filename)
        response['Content-Disposition'] = f"attachment; filename*=UTF-8''{encoded}"

        writer = csv.writer(response)
        writer.writerow(['장비명', '모델', '설치위치', '망구분', '장비유형', 'IP주소', '상태'])
        TYPE_KO = {'switch':'스위치','poe_switch':'PoE스위치','ap':'무선AP','router':'라우터','firewall':'방화벽','server':'서버'}
        STATUS_KO = {'up':'정상','down':'장애','warning':'경고','unknown':'미확인'}
        for d in devices:
            writer.writerow([
                d.name, d.model, d.location, d.network_type,
                TYPE_KO.get(d.device_type, d.device_type),
                d.ip_address or '', STATUS_KO.get(d.status, d.status),
            ])

        # NAS 자동 저장
        try:
            from .services import write_topology_files_to_nas
            write_topology_files_to_nas(school)
        except Exception:
            pass
        return response

    # ── AP위치도 ────────────────────────────────────────
    @action(detail=False, methods=['get'], url_path='ap_map_data')
    def ap_map_data(self, request):
        """AP위치도 — 학교 건물/층별 AP 배치 데이터 + 도면 경로"""
        from apps.schools.models import SchoolEquipment, School
        import os
        school_id = request.query_params.get('school_id')
        if not school_id:
            return Response({'error': 'school_id 필요'}, status=400)
        try:
            school = School.objects.get(pk=school_id)
        except School.DoesNotExist:
            return Response({'error': '학교 없음'}, status=404)

        aps = SchoolEquipment.objects.filter(
            school_id=school_id, category='AP'
        ).order_by('building', 'floor', 'id')

        # 건물/층별 그룹화
        floors = {}
        for ap in aps:
            fkey = f'{ap.building or "본관"}_{ap.floor or "1"}'
            if fkey not in floors:
                bldg = ap.building or '본관'
                flr = ap.floor or '1'
                # 도면 이미지 경로 확인 (NAS)
                plan_dir = f'/app/nas/media/npms/도면/{school.name}'
                plan_file = ''
                for ext in ['png', 'jpg', 'jpeg', 'gif']:
                    candidate = os.path.join(plan_dir, f'{bldg}_{flr}층.{ext}')
                    if os.path.exists(candidate):
                        plan_file = f'/npms/media/npms/도면/{school.name}/{bldg}_{flr}층.{ext}'
                        break
                floors[fkey] = {
                    'key': fkey,
                    'building': bldg,
                    'floor': flr,
                    'plan_image': plan_file,
                    'aps': [],
                }
            floors[fkey]['aps'].append({
                'id': ap.id,
                'model_name': ap.model_name or 'AP',
                'device_id': ap.device_id or '',
                'install_location': ap.install_location or '',
                'network_type': ap.network_type or '',
                'asset_tag': ap.asset_tag or '',
                'pos_x': ap.ap_pos_x,
                'pos_y': ap.ap_pos_y,
            })
        return Response({'school_name': school.name, 'floors': list(floors.values())})

    @action(detail=False, methods=['post'], url_path='ap_map_save')
    def ap_map_save(self, request):
        """AP위치도 좌표 저장"""
        from apps.schools.models import SchoolEquipment
        items = request.data.get('items', [])
        if not items:
            return Response({'error': 'items 필요'}, status=400)
        updated = 0
        for item in items:
            eid = item.get('id')
            if eid:
                SchoolEquipment.objects.filter(pk=eid).update(
                    ap_pos_x=item.get('pos_x'), ap_pos_y=item.get('pos_y'))
                updated += 1
        return Response({'success': True, 'updated': updated})

    @action(detail=False, methods=['post'], url_path='ap_map_upload_plan')
    def ap_map_upload_plan(self, request):
        """AP위치도 — 건물/층 도면 이미지 업로드"""
        from apps.schools.models import School
        import os
        school_id = request.data.get('school_id')
        building = request.data.get('building', '본관')
        floor = request.data.get('floor', '1')
        file = request.FILES.get('file')
        if not school_id or not file:
            return Response({'error': 'school_id와 file 필요'}, status=400)
        try:
            school = School.objects.get(pk=school_id)
        except School.DoesNotExist:
            return Response({'error': '학교 없음'}, status=404)

        plan_dir = f'/app/nas/media/npms/도면/{school.name}'
        os.makedirs(plan_dir, exist_ok=True)
        ext = file.name.rsplit('.', 1)[-1].lower() if '.' in file.name else 'png'
        fname = f'{building}_{floor}층.{ext}'
        fpath = os.path.join(plan_dir, fname)
        with open(fpath, 'wb') as f:
            for chunk in file.chunks():
                f.write(chunk)
        url = f'/npms/media/npms/도면/{school.name}/{fname}'
        return Response({'success': True, 'plan_image': url})

    # ── 랙실장도 ────────────────────────────────────────
    @action(detail=False, methods=['get'], url_path='rack_data')
    def rack_data(self, request):
        """랙실장도 — NAS 파싱 우선, DB 폴백"""
        import os
        from apps.schools.models import School, SchoolEquipment
        school_id = request.query_params.get('school_id')
        if not school_id:
            return Response({'error': 'school_id 필요'}, status=400)
        try:
            school = School.objects.get(pk=school_id)
        except School.DoesNotExist:
            return Response({'error': '학교 없음'}, status=404)

        # NAS 파싱 시도
        nas_result = self._parse_nas_rack(school.name)
        if nas_result:
            return Response(nas_result)

        # DB 폴백
        return Response(self._db_rack(school_id))

    def _parse_nas_rack(self, school_name):
        """NAS 랙실장도 엑셀 파싱"""
        import os
        rack_folder = os.path.join(self.NAS_BASE, self.NAS_FOLDERS.get('rack', ''))
        if not rack_folder or not os.path.isdir(rack_folder):
            return None
        found_file = None
        for fname in os.listdir(rack_folder):
            if school_name in fname and (fname.endswith('.xlsx') or fname.endswith('.xlsm')) and not fname.startswith('~'):
                found_file = os.path.join(rack_folder, fname)
                break
        if not found_file:
            return None

        try:
            import openpyxl
            wb = openpyxl.load_workbook(found_file, read_only=True, data_only=True)
        except Exception:
            return None

        ws = wb[wb.sheetnames[0]]
        max_col = ws.max_column or 13

        # 양식 판별: 첫 U=1 위치 찾기
        u1_row, u1_col = None, None
        for r in range(3, 15):
            for c in [1, 2]:
                v = ws.cell(r, c).value
                if v and str(v).strip() == '1':
                    u1_row, u1_col = r, c
                    break
            if u1_row:
                break

        if not u1_row:
            wb.close()
            return None

        equip_col = u1_col + 1  # 장비명 열

        # 운영장소/랙이름 파싱
        def _get_rack_meta(loc_col, name_col, label_row_offset):
            """운영장소와 랙이름 추출"""
            loc = ''
            name = ''
            # 운영장소: u1_row 위 2~3행
            for r in range(max(1, u1_row - 4), u1_row):
                v = ws.cell(r, loc_col).value
                if v and '서버' in str(v) or v and '전산' in str(v) or v and '층' in str(v):
                    loc = str(v).strip()
                v2 = ws.cell(r, loc_col).value
                if v2 and ('렉' in str(v2) or '랙' in str(v2) or '통신' in str(v2)):
                    name = str(v2).strip()
            return loc, name

        # 랙1 파싱
        racks = []

        def _parse_one_rack(u_col, eq_col, rack_label):
            """한 랙의 장비 목록 파싱"""
            location = ''
            rack_name = ''
            # 운영장소/랙이름: U1 위 행들에서 찾기
            for r in range(max(1, u1_row - 4), u1_row):
                v = ws.cell(r, eq_col).value
                if v:
                    vs = str(v).strip()
                    if '렉' in vs or '랙' in vs or '통신' in vs:
                        rack_name = vs
                    elif len(vs) > 2:
                        location = vs

            items = []
            # 랙 외부 장비 (U1 위 행에서 장비명이 있는 것)
            for r in range(max(1, u1_row - 5), u1_row):
                v = ws.cell(r, eq_col).value
                u_v = ws.cell(r, u_col).value
                if v and not u_v and str(v).strip() not in ('', rack_name, location):
                    vs = str(v).strip()
                    if '통신' not in vs and '실장' not in vs and len(vs) > 1:
                        items.append({'u': 0, 'name': vs, 'type': 'external'})

            # U번호별 장비
            current_u = 0
            for r in range(u1_row, ws.max_row + 1):
                u_v = ws.cell(r, u_col).value
                eq_v = ws.cell(r, eq_col).value
                if u_v and str(u_v).strip().isdigit():
                    current_u = int(u_v)
                if eq_v:
                    name = str(eq_v).strip()
                    if name and name not in (rack_name, location):
                        # 장비 유형 자동 판별
                        etype = 'switch'
                        nl = name.upper()
                        if '패치' in name: etype = 'patch'
                        elif 'FDF' in nl or 'OFD' in nl: etype = 'fdf'
                        elif '방화벽' in name or 'FW' in nl: etype = 'firewall'
                        elif 'UPS' in nl: etype = 'ups'
                        elif '서버' in name: etype = 'server'
                        elif 'P#' in name or 'POE' in nl: etype = 'poe'
                        items.append({'u': current_u, 'name': name, 'type': etype})

            if items:
                racks.append({
                    'rack_name': rack_name or rack_label,
                    'location': location,
                    'items': items,
                    'source': 'nas',
                })

        # 랙1
        _parse_one_rack(u1_col, equip_col, '통신랙 #1')

        # 랙2 확인 (col 8~9 영역)
        if max_col >= 9:
            has_rack2 = False
            for r in range(max(1, u1_row - 3), u1_row + 3):
                for c in [8, 9]:
                    v = ws.cell(r, c).value
                    if v and str(v).strip():
                        has_rack2 = True
                        break
                if has_rack2:
                    break
            if has_rack2:
                # 랙2 U열/장비열 찾기
                u2_col = None
                for c in [7, 8]:
                    v = ws.cell(u1_row, c).value
                    if v and str(v).strip() == '1':
                        u2_col = c
                        break
                if u2_col:
                    _parse_one_rack(u2_col, u2_col + 1, '통신랙 #2')

        wb.close()
        return racks if racks else None

    def _db_rack(self, school_id):
        """DB 기반 랙 데이터 (폴백)"""
        from apps.schools.models import SchoolEquipment
        equips = SchoolEquipment.objects.filter(
            school_id=school_id,
            category__in=['스위치', 'PoE', 'PoE스위치']
        ).order_by('network_type', 'device_id', 'id')
        racks = {}
        for eq in equips:
            rack_key = f'{eq.building or "본관"}_{eq.floor or "1"}_{eq.install_location or "통신실"}'
            if rack_key not in racks:
                racks[rack_key] = {
                    'rack_name': f'통신랙',
                    'location': eq.install_location or '통신실',
                    'items': [],
                    'source': 'db',
                }
            TYPE_COLOR = {'스위치': '#0d6efd', 'PoE': '#6f42c1', 'PoE스위치': '#6f42c1'}
            racks[rack_key]['items'].append({
                'u': eq.rack_unit or 0,
                'name': f'{eq.device_id or ""} {eq.model_name or eq.category}'.strip(),
                'type': 'poe' if eq.category in ('PoE', 'PoE스위치') else 'switch',
                'id': eq.id,
                'rack_unit': eq.rack_unit,
                'rack_size': eq.rack_size or 1,
            })
        return list(racks.values())

    @action(detail=False, methods=['post'], url_path='rack_save')
    def rack_save(self, request):
        """랙실장도 배치 저장 (장비별 U 위치)"""
        from apps.schools.models import SchoolEquipment
        items = request.data.get('items', [])
        if not items:
            return Response({'error': 'items 필요'}, status=400)
        updated = 0
        for item in items:
            eid = item.get('id')
            rack_unit = item.get('rack_unit')
            rack_size = item.get('rack_size', 1)
            if eid:
                SchoolEquipment.objects.filter(pk=eid).update(
                    rack_unit=rack_unit, rack_size=rack_size)
                updated += 1
        return Response({'success': True, 'updated': updated})

    # ── 구성도 (Cytoscape.js) ────────────────────────────
    @action(detail=False, methods=['get'], url_path='diagram_data')
    def diagram_data(self, request):
        """구성도용 장비 데이터 — SchoolEquipment 기반 자동 생성"""
        from apps.schools.models import SchoolEquipment, School
        school_id = request.query_params.get('school_id')
        if not school_id:
            return Response({'error': 'school_id가 필요합니다.'}, status=status.HTTP_400_BAD_REQUEST)
        try:
            school = School.objects.get(pk=school_id)
        except School.DoesNotExist:
            return Response({'error': '학교를 찾을 수 없습니다.'}, status=404)

        # 저장된 배치 데이터가 있으면 반환
        topo = NetworkTopology.objects.filter(school_id=school_id).order_by('-updated_at').first()
        if topo and topo.topology_data.get('diagram_layout'):
            return Response({
                'school_name': school.name,
                'layout': topo.topology_data['diagram_layout'],
                'saved': True,
            })

        # 없으면 SchoolEquipment 기반 자동 생성
        equips = SchoolEquipment.objects.filter(school_id=school_id).order_by('network_type', 'category', 'id')
        # 장비 유형별 분류
        TYPE_MAP = {'스위치': 'switch', 'PoE': 'poe_switch', 'PoE스위치': 'poe_switch',
                    'AP': 'ap', '무선AP': 'ap'}
        nodes = []
        # 방화벽 노드 (인터넷 연결용, 항상 추가)
        nodes.append({
            'id': 'internet', 'label': 'Internet', 'type': 'internet',
            'x': 400, 'y': 30, 'network_type': '',
        })
        nodes.append({
            'id': 'firewall', 'label': '방화벽\n(Secui)', 'type': 'firewall',
            'x': 400, 'y': 120, 'network_type': '',
        })

        # 망구분별 그룹화
        net_groups = {}
        for eq in equips:
            net = eq.network_type or '기타'
            net_groups.setdefault(net, []).append(eq)

        edges = [{'source': 'internet', 'target': 'firewall', 'label': ''}]
        y_offset = 220
        NET_X = {'교사망': 150, '학생망': 400, '무선망': 650}

        for net_name, net_equips in net_groups.items():
            base_x = NET_X.get(net_name, 400)
            # 코어 스위치 (첫 번째 스위치를 코어로)
            core = None
            for eq in net_equips:
                if eq.category in ('스위치', 'PoE', 'PoE스위치'):
                    core = eq
                    break

            if core:
                core_id = f'eq_{core.id}'
                tier_label = f'L{core.tier}' if core.tier else ''
                nodes.append({
                    'id': core_id,
                    'label': f'{core.model_name or core.category}\n{core.device_id or ""}\n{tier_label}'.strip(),
                    'type': TYPE_MAP.get(core.category, 'switch'),
                    'x': base_x, 'y': y_offset,
                    'network_type': net_name,
                    'equip_id': core.id,
                })
                edges.append({'source': 'firewall', 'target': core_id, 'label': net_name})

                # 나머지 장비는 코어에 연결
                sub_y = y_offset + 110
                sub_x = base_x - 80
                for i, eq in enumerate(net_equips):
                    if eq.id == core.id:
                        continue
                    eq_id = f'eq_{eq.id}'
                    nodes.append({
                        'id': eq_id,
                        'label': f'{eq.model_name or eq.category}\n{eq.device_id or ""}'.strip(),
                        'type': TYPE_MAP.get(eq.category, 'switch'),
                        'x': sub_x + (i % 4) * 120,
                        'y': sub_y + (i // 4) * 100,
                        'network_type': net_name,
                        'equip_id': eq.id,
                    })
                    edges.append({'source': core_id, 'target': eq_id, 'label': ''})

        return Response({
            'school_name': school.name,
            'layout': {'nodes': nodes, 'edges': edges},
            'saved': False,
        })

    @action(detail=False, methods=['post'], url_path='diagram_save')
    def diagram_save(self, request):
        """구성도 배치 저장 (Cytoscape.js에서 드래그한 위치)"""
        from apps.schools.models import School
        school_id = request.data.get('school_id')
        layout = request.data.get('layout', {})
        if not school_id or not layout:
            return Response({'error': 'school_id와 layout이 필요합니다.'}, status=400)
        try:
            school = School.objects.get(pk=school_id)
        except School.DoesNotExist:
            return Response({'error': '학교를 찾을 수 없습니다.'}, status=404)

        topo, _ = NetworkTopology.objects.get_or_create(
            school=school,
            defaults={'topology_data': {}}
        )
        data = topo.topology_data or {}
        data['diagram_layout'] = layout
        topo.topology_data = data
        topo.save(update_fields=['topology_data', 'updated_at'])
        return Response({'success': True})

    @action(detail=False, methods=['get'], url_path='diagram_pptx')
    def diagram_pptx(self, request):
        """구성도 PPTX 다운로드"""
        import io
        from urllib.parse import quote
        from django.http import HttpResponse
        from apps.schools.models import School
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt, Emu
            from pptx.dml.color import RGBColor
            from pptx.enum.text import PP_ALIGN
        except ImportError:
            return HttpResponse('python-pptx 필요', status=500)

        school_id = request.query_params.get('school_id')
        if not school_id:
            return HttpResponse('school_id 필요', status=400)
        try:
            school = School.objects.get(pk=school_id)
        except School.DoesNotExist:
            return HttpResponse('학교 없음', status=404)

        # 구성도 데이터 로드
        topo = NetworkTopology.objects.filter(school_id=school_id).order_by('-updated_at').first()
        layout = None
        if topo and topo.topology_data.get('diagram_layout'):
            layout = topo.topology_data['diagram_layout']

        if not layout:
            # 자동 생성 데이터 사용
            from django.test import RequestFactory
            factory = RequestFactory()
            fake_req = factory.get(f'/?school_id={school_id}')
            fake_req.user = request.user
            resp = self.diagram_data(fake_req)
            layout = resp.data.get('layout', {})

        nodes = layout.get('nodes', [])
        edges = layout.get('edges', [])

        prs = Presentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # 빈 슬라이드

        # 제목
        txBox = slide.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(5), Inches(0.5))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = f'{school.name} — 네트워크 구성도'
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0x1F, 0x49, 0x7D)

        # 노드 색상
        TYPE_COLORS = {
            'internet':   RGBColor(0x6C, 0x75, 0x7D),
            'firewall':   RGBColor(0xDC, 0x35, 0x45),
            'switch':     RGBColor(0x0D, 0x6E, 0xFD),
            'poe_switch': RGBColor(0x6F, 0x42, 0xC1),
            'ap':         RGBColor(0x19, 0x87, 0x54),
        }

        # 좌표 스케일 (Cytoscape px → PPTX Inches)
        if nodes:
            max_x = max(n.get('x', 0) for n in nodes) or 800
            max_y = max(n.get('y', 0) for n in nodes) or 600
        else:
            max_x, max_y = 800, 600
        scale_x = 11.0 / max(max_x, 1)
        scale_y = 5.5 / max(max_y, 1)

        node_positions = {}
        for n in nodes:
            nx = n.get('x', 0) * scale_x + 1.0
            ny = n.get('y', 0) * scale_y + 1.2
            w, h = 1.3, 0.7
            node_positions[n['id']] = (nx + w/2, ny + h/2)

            color = TYPE_COLORS.get(n.get('type', ''), RGBColor(0x6C, 0x75, 0x7D))
            shape = slide.shapes.add_shape(
                1,  # MSO_SHAPE.ROUNDED_RECTANGLE
                Emu(int(nx * 914400)), Emu(int(ny * 914400)),
                Emu(int(w * 914400)), Emu(int(h * 914400))
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = color
            shape.line.color.rgb = color
            tf = shape.text_frame
            tf.word_wrap = True
            for line in (n.get('label', '') or '').split('\n'):
                p = tf.add_paragraph() if tf.paragraphs[0].text else tf.paragraphs[0]
                p.text = line.strip()
                p.font.size = Pt(8)
                p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                p.alignment = PP_ALIGN.CENTER

        # 엣지 (직선 커넥터)
        for e in edges:
            src = node_positions.get(e.get('source'))
            tgt = node_positions.get(e.get('target'))
            if not src or not tgt:
                continue
            connector = slide.shapes.add_connector(
                1,  # MSO_CONNECTOR.STRAIGHT
                Emu(int(src[0] * 914400)), Emu(int(src[1] * 914400)),
                Emu(int(tgt[0] * 914400)), Emu(int(tgt[1] * 914400))
            )
            connector.line.color.rgb = RGBColor(0x99, 0x99, 0x99)
            connector.line.width = Pt(1.5)

        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)
        fname = f'구성도_{school.name}.pptx'
        resp = HttpResponse(buf.read(),
            content_type='application/vnd.openxmlformats-officedocument.presentationml.presentation')
        resp['Content-Disposition'] = f"attachment; filename*=UTF-8''{quote(fname)}"
        return resp

    # ── 선번장 ────────────────────────────────────────
    @action(detail=False, methods=['get'], url_path='portmap_data')
    def portmap_data(self, request):
        """선번장 — DB 우선 (사전 파싱 캐시), NAS 실시간 파싱 폴백"""
        import os
        from apps.schools.models import School, SchoolEquipment
        school_id = request.query_params.get('school_id')
        if not school_id:
            return Response({'error': 'school_id 필요'}, status=400)
        try:
            school = School.objects.get(pk=school_id)
        except School.DoesNotExist:
            return Response({'error': '학교 없음'}, status=404)

        # 1순위: DB에 port_map 데이터가 실제로 있으면 즉시 반환
        has_portmap = SchoolEquipment.objects.filter(
            school=school, category__in=['스위치', 'PoE', 'PoE스위치'],
            port_map__isnull=False,
        ).exclude(port_map=[]).exists()
        if has_portmap:
            return Response(self._db_portmap(school_id))

        # 2순위: NAS 실시간 파싱
        nas_result = self._parse_nas_portmap(school.name)
        if nas_result:
            return Response(nas_result)

        # 3순위: DB 장비만 있고 port_map 없는 경우 (빈 포트)
        db_result = self._db_portmap(school_id)
        if db_result:
            return Response(db_result)

        return Response([])

    @staticmethod
    def _detect_cable_by_color(ws, row, col, cable_labels=None):
        """셀 배경색으로 실제 연결 케이블 감지
        각 포트에 3칸(C6|C5.e|C5 또는 UTP|SM|MM)이 있고,
        배경색이 칠해진(patternType=solid) 칸이 실제 연결된 케이블.
        theme=0(흰색) 배경은 제외.
        """
        if cable_labels is None:
            cable_labels = ['C6', 'C5.e', 'C5']
        for offset, label in enumerate(cable_labels):
            c = ws.cell(row, col + offset)
            if c.fill and c.fill.patternType == 'solid':
                fg = c.fill.fgColor
                if not fg:
                    continue
                # theme 기반 색상(theme=0은 흰색) 제외
                try:
                    theme = fg.theme
                    if theme is not None and isinstance(theme, int):
                        continue  # theme 색상은 실제 케이블 색상이 아님
                except (TypeError, AttributeError):
                    pass
                # RGB 명시 색상만 인정
                try:
                    rgb = fg.rgb
                    if rgb and isinstance(rgb, str) and rgb.startswith('FF') and rgb != 'FF000000':
                        return label
                except (TypeError, AttributeError):
                    continue
        return ''

    def _parse_nas_portmap(self, school_name):
        """NAS 선번장 파일 파싱 — xlsx/xlsm 양식 A/B 자동 감지"""
        import os
        folder = os.path.join(self.NAS_BASE, self.NAS_FOLDERS['portmap'])
        if not os.path.isdir(folder):
            return None
        found_file = None
        for fname in os.listdir(folder):
            if school_name in fname:
                found_file = os.path.join(folder, fname)
                break
        if not found_file:
            return None

        try:
            import openpyxl
            wb = openpyxl.load_workbook(found_file, data_only=True)
        except Exception:
            return None

        ext = found_file.rsplit('.', 1)[-1].lower()
        result = []

        for sname in wb.sheetnames:
            if sname in ('장비스펙', 'FDF', 'FDF현황'):
                continue
            ws = wb[sname]
            if not ws.max_row:
                continue

            # 양식 자동 감지: xlsm(67열,기준행C3) vs xlsx(66열,기준행A1)
            is_type_b = ext == 'xlsm'
            if not is_type_b:
                # xlsx도 양식B일 수 있음 — C3에 '허' 가 있는지 확인
                for r in range(1, min(10, ws.max_row + 1)):
                    v = ws.cell(r, 3).value
                    if v and '허' in str(v):
                        is_type_b = True
                        break

            if is_type_b:
                switches = self._parse_type_b(ws, sname)
            else:
                switches = self._parse_type_a(ws, sname)
            result.extend(switches)

        wb.close()
        return result if result else None

    def _parse_type_b(self, ws, sname):
        """양식 B (xlsm 스타일): 67열, 19행 간격, C열 시작
        케이블: 셀 배경색으로 감지 (C6|C5.e|C5 3칸, 포트25~는 UTP|SM|MM)
        """
        PORT_COLS_B = [3, 7, 11, 15, 19, 23, 27, 31, 35, 39, 43, 47, 53, 57]
        UPLINK_COLS_B = [53, 57]  # 포트 25~28 (UTP|SM|MM)
        switches = []
        for r in range(1, ws.max_row + 1):
            v = ws.cell(r, 3).value
            if not (v and '허' in str(v) and 'I' in str(v)):
                continue
            base = r
            hub_id = str(ws.cell(base, 7).value or '').strip()
            model = str(ws.cell(base, 19).value or '').strip()
            location = str(ws.cell(base, 31).value or '').strip()
            net_type = str(ws.cell(base + 2, 7).value or '').strip()
            mfr = str(ws.cell(base + 2, 31).value or '').strip()
            if not hub_id:
                continue

            ports = []
            for col in PORT_COLS_B:
                labels = ['UTP', 'SM', 'MM'] if col in UPLINK_COLS_B else ['C6', 'C5.e', 'C5']
                # 홀수
                pnum = ws.cell(base + 6, col).value
                conn = str(ws.cell(base + 8, col).value or '').strip()
                cable = self._detect_cable_by_color(ws, base + 7, col, labels)
                if pnum and str(pnum).strip():
                    pn = int(pnum) if str(pnum).strip().isdigit() else 0
                    if pn > 0:
                        ports.append({'port': pn, 'connected_to': conn if conn != '/' else '',
                                      'cable': cable, 'vlan': '', 'note': ''})
                # 짝수
                pnum2 = ws.cell(base + 11, col).value
                conn2 = str(ws.cell(base + 13, col).value or '').strip()
                cable2 = self._detect_cable_by_color(ws, base + 12, col, labels)
                if pnum2 and str(pnum2).strip():
                    pn2 = int(pnum2) if str(pnum2).strip().isdigit() else 0
                    if pn2 > 0:
                        ports.append({'port': pn2, 'connected_to': conn2 if conn2 != '/' else '',
                                      'cable': cable2, 'vlan': '', 'note': ''})
            ports.sort(key=lambda p: p['port'])
            switches.append({
                'id': 0, 'model_name': model, 'device_id': hub_id,
                'category': 'PoE' if 'P' in hub_id.upper() and '#' in hub_id else '스위치',
                'network_type': net_type, 'manufacturer': mfr,
                'building': '', 'floor': '', 'install_location': location,
                'port_count': len(ports), 'ports': ports, 'sheet': sname,
                'source': 'nas',
            })
        return switches

    def _parse_type_a(self, ws, sname):
        """양식 A (xlsx 스타일): 66열, 10행 간격, A열 시작
        케이블: 셀 배경색으로 감지 (C6|C5.e|C5 3칸, 업링크는 UTP|SM|MM)
        """
        PORT_COLS_A = [1, 4, 7, 10, 13, 16, 19, 22, 25, 28, 31, 34, 37, 40, 43, 46, 49, 52, 55, 58]
        UPLINK_COLS_A = [55, 58]  # 포트 25~28 (UTP|SM|MM)
        switches = []
        for r in range(1, ws.max_row + 1):
            v = ws.cell(r, 1).value
            if not (v and str(v).strip() == '망구분'):
                continue
            base = r
            net_type = str(ws.cell(base, 4).value or '').strip()
            model = str(ws.cell(base, 13).value or '').strip()
            location = str(ws.cell(base, 22).value or '').strip()
            poe = str(ws.cell(base, 31).value or '').strip()
            hub_id = str(ws.cell(base + 1, 4).value or '').strip()  # ID 행
            mfr = str(ws.cell(base + 1, 13).value or '').strip()
            speed = str(ws.cell(base + 1, 31).value or '').strip()
            if not hub_id and not model:
                continue

            ports = []
            for col in PORT_COLS_A:
                labels = ['UTP', 'SM', 'MM'] if col in UPLINK_COLS_A else ['C6', 'C5.e', 'C5']
                # 홀수 (base+2)
                pnum = ws.cell(base + 2, col).value
                cable = self._detect_cable_by_color(ws, base + 3, col, labels)
                conn = str(ws.cell(base + 4, col).value or '').strip()
                if pnum and str(pnum).strip().isdigit():
                    pn = int(pnum)
                    ports.append({'port': pn, 'connected_to': conn if conn and conn != '/' else '',
                                  'cable': cable, 'vlan': '', 'note': ''})
                # 짝수 (base+5)
                pnum2 = ws.cell(base + 5, col).value
                cable2 = self._detect_cable_by_color(ws, base + 6, col, labels)
                conn2 = str(ws.cell(base + 7, col).value or '').strip()
                if pnum2 and str(pnum2).strip().isdigit():
                    pn2 = int(pnum2)
                    ports.append({'port': pn2, 'connected_to': conn2 if conn2 and conn2 != '/' else '',
                                  'cable': cable2, 'vlan': '', 'note': ''})
            ports.sort(key=lambda p: p['port'])
            switches.append({
                'id': 0, 'model_name': model, 'device_id': hub_id or model,
                'category': 'PoE' if poe and poe not in ('N', '없음', '') else '스위치',
                'network_type': net_type, 'manufacturer': mfr,
                'building': '', 'floor': '', 'install_location': location,
                'port_count': len(ports), 'ports': ports, 'sheet': sname,
                'source': 'nas',
            })
        return switches

    def _db_portmap(self, school_id):
        """DB 기반 포트맵 (폴백)"""
        from apps.schools.models import SchoolEquipment
        equips = SchoolEquipment.objects.filter(
            school_id=school_id, category__in=['스위치', 'PoE', 'PoE스위치']
        ).order_by('network_type', 'device_id', 'id')
        result = []
        for eq in equips:
            port_count = 24
            model = (eq.model_name or '').upper()
            if '48' in model: port_count = 48
            elif '16' in model: port_count = 16
            elif '8' in model and '28' not in model: port_count = 8
            elif '28' in model: port_count = 28
            elif '52' in model: port_count = 52
            port_map = eq.port_map or []
            existing = {p.get('port'): p for p in port_map if isinstance(p, dict)}
            ports = []
            for i in range(1, port_count + 1):
                p = existing.get(i, {})
                ports.append({'port': i, 'connected_to': p.get('connected_to', ''),
                              'vlan': p.get('vlan', ''), 'cable': p.get('cable', ''),
                              'note': p.get('note', ''), 'status': p.get('status', 'down')})
            result.append({
                'id': eq.id, 'model_name': eq.model_name or '', 'device_id': eq.device_id or '',
                'category': eq.category, 'network_type': eq.network_type or '',
                'building': eq.building or '', 'floor': eq.floor or '',
                'install_location': eq.install_location or '',
                'port_count': port_count, 'ports': ports, 'source': 'db',
            })
        return result

    @action(detail=False, methods=['post'], url_path='portmap_save')
    def portmap_save(self, request):
        """선번장 포트 연결정보 저장"""
        from apps.schools.models import SchoolEquipment
        equip_id = request.data.get('equipment_id')
        ports = request.data.get('ports', [])
        if not equip_id:
            return Response({'error': 'equipment_id 필요'}, status=400)
        try:
            eq = SchoolEquipment.objects.get(pk=equip_id)
        except SchoolEquipment.DoesNotExist:
            return Response({'error': '장비 없음'}, status=404)
        eq.port_map = ports
        eq.save(update_fields=['port_map'])
        return Response({'success': True})

    @action(detail=False, methods=['get'], url_path='portmap_pdf')
    def portmap_pdf(self, request):
        """선번장 PDF 다운로드 — 가로 방향, 스위치별 1페이지"""
        from urllib.parse import quote
        from django.http import HttpResponse
        from apps.schools.models import School, SchoolEquipment
        import weasyprint

        school_id = request.query_params.get('school_id')
        if not school_id:
            return HttpResponse('school_id 필요', status=400)
        try:
            school = School.objects.get(pk=school_id)
        except School.DoesNotExist:
            return HttpResponse('학교 없음', status=404)

        # 데이터 수집: DB port_map 우선, 없으면 NAS 파싱
        has_portmap = SchoolEquipment.objects.filter(
            school=school, category__in=['스위치', 'PoE', 'PoE스위치'],
            port_map__isnull=False,
        ).exclude(port_map=[]).exists()

        if has_portmap:
            switches = self._db_portmap(school_id)
        else:
            switches = self._parse_nas_portmap(school.name) or self._db_portmap(school_id)

        if not switches:
            return HttpResponse('선번장 데이터 없음', status=404)

        # 케이블 색상 매핑
        cable_colors = {
            'C6': {'bg': '#198754', 'text': '#fff', 'label': 'Cat6'},
            'C5.e': {'bg': '#ffc107', 'text': '#333', 'label': 'Cat5E'},
            'C5': {'bg': '#e63946', 'text': '#fff', 'label': 'Cat5'},
            'UTP': {'bg': '#6f42c1', 'text': '#fff', 'label': 'UTP'},
            'SM': {'bg': '#0dcaf0', 'text': '#333', 'label': 'SM'},
            'MM': {'bg': '#ff6b00', 'text': '#fff', 'label': 'MM'},
        }
        def get_cable_style(cable):
            c = (cable or '').upper().replace('.', '').replace(' ', '')
            if 'C6' in c or 'CAT6' in c: return cable_colors['C6']
            if 'C5E' in c or 'CAT5E' in c or '5E' in c: return cable_colors['C5.e']
            if 'C5' in c or 'CAT5' in c: return cable_colors['C5']
            if 'UTP' in c: return cable_colors['UTP']
            if 'SM' in c: return cable_colors['SM']
            if 'MM' in c: return cable_colors['MM']
            return None

        # HTML 생성
        pages_html = ''
        for idx, sw in enumerate(switches):
            ports = sw.get('ports', [])
            pc = sw.get('port_count', len(ports)) or len(ports)
            conn_count = sum(1 for p in ports if p.get('cable', '').strip())
            poe = 'Y' if sw.get('category') in ('PoE', 'PoE스위치') or 'P#' in (sw.get('device_id') or '').upper() else 'N'
            location = ' / '.join(filter(None, [sw.get('building'), sw.get('floor'), sw.get('install_location')])) or '-'

            # SVG 포트맵
            # 포트 status 기본값: 케이블 있으면 활성, 없으면 미연결
            for p in ports:
                if not p.get('status'):
                    p['status'] = 'up' if p.get('cable', '').strip() else 'none'

            odd_ports = [p for p in ports if p.get('port', 0) % 2 == 1]
            even_ports = [p for p in ports if p.get('port', 0) % 2 == 0]
            half_cols = max(len(odd_ports), len(even_ports), 1)
            pw, ph, gap, px, py = 28, 22, 2, 20, 18
            svg_w = max(px * 2 + half_cols * (pw + gap), 220)
            led_r, led_gap = 3, 5
            row_h = ph + led_r * 2 + led_gap + 10
            svg_h = py + row_h * 2 + 18

            def svg_row(port_list, y):
                s = ''
                for i, p in enumerate(port_list):
                    x = px + i * (pw + gap)
                    cs = get_cable_style(p.get('cable'))
                    fill = cs['bg'] if cs else '#fff'
                    border = '#b0bec5' if not cs else fill
                    txt_c = cs['text'] if cs else '#546e7a'
                    s += f'<rect x="{x}" y="{y}" width="{pw}" height="{ph}" rx="2" fill="{fill}" stroke="{border}" stroke-width="1"/>'
                    s += f'<text x="{x+pw//2}" y="{y+ph//2+4}" fill="{txt_c}" font-size="7" text-anchor="middle" font-weight="bold">{p.get("port","")}</text>'
                    # LED
                    led_y = y + ph + led_gap
                    st = p.get('status', '')
                    led_c = '#4caf50' if st == 'up' else '#f44336' if st == 'down' else '#bdbdbd'
                    s += f'<circle cx="{x+pw//2}" cy="{led_y}" r="{led_r}" fill="{led_c}" stroke="#fff" stroke-width=".5"/>'
                return s

            svg = f'<svg width="{svg_w}" height="{svg_h}" xmlns="http://www.w3.org/2000/svg">'
            svg += f'<rect width="{svg_w}" height="{svg_h}" rx="6" fill="#e8edf2"/>'
            svg += f'<rect x="4" y="4" width="{svg_w-8}" height="{svg_h-8}" rx="4" fill="#f0f4f8" stroke="#b0bec5" stroke-width="1"/>'
            svg += f'<text x="{svg_w//2}" y="{py-4}" fill="#455a64" font-size="7" text-anchor="middle" font-weight="600">{sw.get("device_id","")} {sw.get("model_name","")}</text>'
            svg += svg_row(odd_ports, py)
            svg += svg_row(even_ports, py + row_h)
            # 범례 1줄: 좌측 케이블 + 우측 LED
            ly = py + row_h * 2 + 8
            legend = ''
            for key, cs in cable_colors.items():
                legend += f'<tspan fill="{cs["bg"]}">&#9632;</tspan>{cs["label"]} '
            legend += '<tspan fill="#ccc">&#9633;</tspan>미연결'
            svg += f'<text x="{px}" y="{ly}" fill="#607d8b" font-size="6">{legend}</text>'
            svg += f'<text x="{svg_w - 110}" y="{ly}" fill="#607d8b" font-size="6">'
            svg += '<tspan fill="#4caf50">&#9679;</tspan>활성 <tspan fill="#f44336">&#9679;</tspan>비활성 <tspan fill="#bdbdbd">&#9679;</tspan>미연결'
            svg += '</text>'
            svg += '</svg>'

            # 포트 표
            port_rows = ''
            for p in ports:
                cs = get_cable_style(p.get('cable'))
                cable_display = cs['label'] if cs else '-'
                cable_dot = f'<span style="display:inline-block;width:6px;height:6px;border-radius:50%;background:{cs["bg"]};margin-right:3px;"></span>' if cs else ''
                status = p.get('status', '')
                status_label = '활성' if status == 'up' else '비활성' if status == 'down' else '미연결'
                status_color = '#198754' if status == 'up' else '#dc3545' if status == 'down' else '#999'
                port_rows += f'''<tr>
                    <td style="text-align:center;font-weight:700;padding:2px 4px;">{p.get('port','')}</td>
                    <td style="padding:2px 6px;">{p.get('connected_to','') or ''}</td>
                    <td style="text-align:center;padding:2px 4px;">{cable_dot}{cable_display}</td>
                    <td style="text-align:center;padding:2px 4px;color:{status_color};font-weight:600;">{status_label}</td>
                </tr>'''

            page_break = 'page-break-before:always;' if idx > 0 else ''
            pages_html += f'''
            <div class="page" style="{page_break}">
                <div class="header">
                    <span class="school-name">{school.name}</span>
                    <span class="title">네트워크 선번장</span>
                    <span class="page-num">{idx+1} / {len(switches)}</span>
                </div>
                <div class="content">
                    <div class="left-panel">
                        <table class="info-table">
                            <tr><th>스위치ID</th><td class="bold">{sw.get('device_id','') or '-'}</td></tr>
                            <tr><th>망구분</th><td>{sw.get('network_type','') or '-'}</td></tr>
                            <tr><th>모델명</th><td>{sw.get('model_name','') or '-'}</td></tr>
                            <tr><th>제조사</th><td>{sw.get('manufacturer','') or '-'}</td></tr>
                            <tr><th>운영장소</th><td>{location}</td></tr>
                            <tr><th>PoE</th><td>{'<span class="badge-y">Y</span>' if poe=='Y' else '<span class="badge-n">N</span>'}</td></tr>
                            <tr><th>포트현황</th><td><span class="badge-conn">{conn_count}</span> / {pc} 연결</td></tr>
                        </table>
                        <div class="svg-box">{svg}</div>
                    </div>
                    <div class="right-panel">
                        <table class="port-table" style="table-layout:fixed;width:100%;">
                            <colgroup><col style="width:32px;"><col style="width:40%;"><col style="width:30%;"><col style="width:30%;"></colgroup>
                            <thead><tr>
                                <th>Port</th>
                                <th>연결장비</th>
                                <th>케이블</th>
                                <th>점검상태</th>
                            </tr></thead>
                            <tbody>{port_rows}</tbody>
                        </table>
                    </div>
                </div>
            </div>'''

        html = f'''<!DOCTYPE html>
<html><head><meta charset="utf-8">
<style>
@page {{ size: A4 landscape; margin: 12mm 10mm; }}
body {{ font-family: 'Noto Sans KR', 'Malgun Gothic', sans-serif; font-size: 8pt; color: #333; margin: 0; }}
.page {{ width: 100%; }}
.header {{ border-bottom: 2px solid #1565c0; padding-bottom: 2px; margin-bottom: 4px; overflow: hidden; }}
.school-name {{ font-size: 12pt; font-weight: 800; color: #1565c0; }}
.title {{ font-size: 10pt; font-weight: 600; color: #455a64; margin-left: 12px; }}
.page-num {{ float: right; font-size: 8pt; color: #90a4ae; }}
.content {{ overflow: hidden; }}
.left-panel {{ float: left; width: 46%; padding-right: 10px; }}
.right-panel {{ float: right; width: 52%; }}
.info-table {{ width: 100%; border-collapse: collapse; margin-bottom: 8px; font-size: 7.5pt; }}
.info-table {{ border-collapse: collapse; }}
.info-table th {{ background: #f5f7fa; color: #546e7a; font-weight: 600; font-size: 6.5pt; text-align: left; padding: 0 6px; border: 0.5px solid #dee2e6; width: 60px; height: 5mm; line-height: 5mm; }}
.info-table td {{ padding: 0 6px; border: 0.5px solid #dee2e6; height: 5mm; line-height: 5mm; font-size: 6.5pt; }}
.info-table td.bold {{ font-weight: 700; color: #1565c0; }}
.badge-y {{ background: #1565c0; color: #fff; padding: 1px 6px; border-radius: 3px; font-size: 7pt; font-weight: 700; }}
.badge-n {{ background: #90a4ae; color: #fff; padding: 1px 6px; border-radius: 3px; font-size: 7pt; }}
.badge-conn {{ background: #198754; color: #fff; padding: 1px 6px; border-radius: 3px; font-size: 7pt; font-weight: 700; }}
.svg-box {{ text-align: center; margin-top: 6px; }}
.svg-box svg {{ width: 100%; height: auto; }}
.port-table {{ width: 100%; border-collapse: collapse; font-size: 7pt; }}
.port-table {{ border-collapse: collapse; }}
.port-table thead th {{ background: #37474f; color: #eceff1; font-weight: 600; font-size: 6.5pt; padding: 0 3px; text-align: center; height: 5mm; line-height: 5mm; }}
.port-table tbody tr:nth-child(even) {{ background: #f8fafc; }}
.port-table tbody tr:nth-child(odd) {{ background: #fff; }}
.port-table tbody td {{ border-bottom: 0.5px solid #eceff1; font-size: 6.5pt; height: 5mm; line-height: 5mm; padding: 0 3px; }}
</style></head><body>{pages_html}</body></html>'''

        pdf = weasyprint.HTML(string=html).write_pdf()
        filename = f'선번장_{school.name}.pdf'
        response = HttpResponse(pdf, content_type='application/pdf')
        response['Content-Disposition'] = f"attachment; filename*=UTF-8''{quote(filename)}"
        return response

    @action(detail=False, methods=['get'], url_path='portmap_excel')
    def portmap_excel(self, request):
        """선번장 엑셀 다운로드"""
        import io
        from urllib.parse import quote
        from django.http import HttpResponse
        from apps.schools.models import School, SchoolEquipment
        try:
            import openpyxl
            from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
            from openpyxl.utils import get_column_letter
        except ImportError:
            return HttpResponse('openpyxl 필요', status=500)

        school_id = request.query_params.get('school_id')
        if not school_id:
            return HttpResponse('school_id 필요', status=400)
        try:
            school = School.objects.get(pk=school_id)
        except School.DoesNotExist:
            return HttpResponse('학교 없음', status=404)

        equips = SchoolEquipment.objects.filter(
            school=school, category__in=['스위치', 'PoE', 'PoE스위치']
        ).order_by('network_type', 'device_id', 'id')

        wb = openpyxl.Workbook()
        hdr_fill = PatternFill('solid', fgColor='1F497D')
        hdr_font = Font(bold=True, color='FFFFFF', size=10)
        ctr = Alignment(horizontal='center', vertical='center')
        thin = Border(left=Side('thin'), right=Side('thin'),
                      top=Side('thin'), bottom=Side('thin'))

        first = True
        for eq in equips:
            label = f'{eq.device_id or eq.model_name or eq.category}'[:31]
            if first:
                ws = wb.active
                ws.title = label
                first = False
            else:
                ws = wb.create_sheet(label)

            ws.merge_cells('A1:G1')
            title = f'{school.name} — {eq.model_name or eq.category} ({eq.device_id or ""}) [{eq.network_type or ""}]'
            ws.cell(1, 1, title).font = Font(bold=True, size=12)
            ws.cell(2, 1, f'위치: {eq.building or ""} {eq.floor or ""}층 {eq.install_location or ""}').font = Font(size=9, color='666666')

            headers = ['포트', '연결 장비', 'VLAN', '케이블', '상태', '비고']
            for ci, h in enumerate(headers, 1):
                c = ws.cell(4, ci, h)
                c.font, c.fill, c.alignment, c.border = hdr_font, hdr_fill, ctr, thin

            port_map = eq.port_map or []
            model = (eq.model_name or '').upper()
            port_count = 48 if '48' in model else 52 if '52' in model else 28 if '28' in model else 16 if '16' in model else 8 if ('8' in model and '28' not in model) else 24
            existing = {p.get('port'): p for p in port_map if isinstance(p, dict)}

            STATUS_KO = {'up': '활성', 'down': '비활성', 'disabled': '비사용'}
            for i in range(1, port_count + 1):
                p = existing.get(i, {})
                vals = [i, p.get('connected_to', ''), p.get('vlan', ''),
                        p.get('cable', ''), STATUS_KO.get(p.get('status', 'down'), p.get('status', '')),
                        p.get('note', '')]
                for ci, v in enumerate(vals, 1):
                    c = ws.cell(i + 4, ci, v)
                    c.border = thin
                    if ci == 1:
                        c.alignment = ctr

            col_widths = [8, 25, 10, 12, 10, 25]
            for ci, w in enumerate(col_widths, 1):
                ws.column_dimensions[get_column_letter(ci)].width = w

        if first:  # 스위치 장비가 없는 경우
            ws = wb.active
            ws.cell(1, 1, '스위치/PoE 장비가 없습니다.')

        buf = io.BytesIO()
        wb.save(buf)
        buf.seek(0)
        fname = f'선번장_{school.name}.xlsx'
        resp = HttpResponse(buf.read(),
            content_type='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet')
        resp['Content-Disposition'] = f"attachment; filename*=UTF-8''{quote(fname)}"
        return resp

    # ── NAS 파일 기반 뷰어 ────────────────────────────
    NAS_BASE = '/app/nas/media/npms/산출물/2025년 테크센터'
    NAS_FOLDERS = {
        'diagram':  '2025년 테크센터-네트워크 구성도',
        'portmap':  '2025년 테크센터-네트워크 선번장',
        'rack':     '2025년 테크센터-네트워크 통신랙실장도',
        'apmap':    '2025년 테크센터-건물 정보',
    }
    NAS_PREFIXES = {
        'diagram':  '2025년 테크센터-네트워크 구성도_',
        'portmap':  '2025년 테크센터-네트워크 선번장_',
        'rack':     '2025년 테크센터-네트워크 통신랙실장도_',
        'apmap':    '2025년 테크센터-건물정보_',
    }

    @action(detail=False, methods=['get'], url_path='nas_file_info')
    def nas_file_info(self, request):
        """NAS 파일 슬라이드/시트/페이지 정보 조회"""
        import os
        from apps.schools.models import School
        school_id = request.query_params.get('school_id')
        tab = request.query_params.get('tab')  # diagram|portmap|rack|apmap
        if not school_id or not tab:
            return Response({'error': 'school_id와 tab 필요'}, status=400)
        try:
            school = School.objects.get(pk=school_id)
        except School.DoesNotExist:
            return Response({'error': '학교 없음'}, status=404)

        folder = self.NAS_FOLDERS.get(tab)
        prefix = self.NAS_PREFIXES.get(tab)
        if not folder or not prefix:
            return Response({'error': '잘못된 tab'}, status=400)

        nas_dir = os.path.join(self.NAS_BASE, folder)
        # 학교명으로 파일 찾기
        school_name = school.name
        found_file = None
        if os.path.isdir(nas_dir):
            for fname in os.listdir(nas_dir):
                if school_name in fname:
                    found_file = os.path.join(nas_dir, fname)
                    break

        if not found_file or not os.path.isfile(found_file):
            return Response({'found': False, 'school_name': school_name, 'tab': tab})

        ext = found_file.rsplit('.', 1)[-1].lower()
        tabs_info = []

        import subprocess, tempfile, shutil, hashlib
        try:
            if ext == 'pdf':
                import fitz
                doc = fitz.open(found_file)
                for i in range(len(doc)):
                    tabs_info.append({'index': i, 'name': f'페이지 {i+1}'})
                doc.close()
            else:
                # PPTX/Excel → LibreOffice로 PDF 변환 후 페이지 수 확인
                file_hash = hashlib.md5(found_file.encode() + str(os.path.getmtime(found_file)).encode()).hexdigest()[:12]
                cache_dir = f'/tmp/nms_cache/{file_hash}'
                pdf_path = os.path.join(cache_dir, 'converted.pdf')
                if not os.path.exists(pdf_path):
                    os.makedirs(cache_dir, exist_ok=True)
                    with tempfile.TemporaryDirectory() as tmpdir:
                        subprocess.run(
                            ['libreoffice', '--headless', '--convert-to', 'pdf', '--outdir', tmpdir, found_file],
                            timeout=120, capture_output=True)
                        for f in os.listdir(tmpdir):
                            if f.endswith('.pdf'):
                                shutil.copy2(os.path.join(tmpdir, f), pdf_path)
                                break
                if os.path.exists(pdf_path):
                    import fitz
                    doc = fitz.open(pdf_path)
                    # PPTX는 슬라이드명, Excel은 시트명 추출
                    if ext == 'pptx':
                        try:
                            from pptx import Presentation
                            prs = Presentation(found_file)
                            for i, slide in enumerate(prs.slides):
                                title = slide.shapes.title.text if slide.shapes.title else ''
                                tabs_info.append({'index': i, 'name': title or f'슬라이드 {i+1}'})
                        except Exception:
                            for i in range(len(doc)):
                                tabs_info.append({'index': i, 'name': f'슬라이드 {i+1}'})
                    elif ext in ('xlsx', 'xlsm'):
                        try:
                            import openpyxl
                            wb = openpyxl.load_workbook(found_file, read_only=True, data_only=True)
                            for i, sname in enumerate(wb.sheetnames):
                                tabs_info.append({'index': i, 'name': sname})
                            wb.close()
                        except Exception:
                            for i in range(len(doc)):
                                tabs_info.append({'index': i, 'name': f'시트 {i+1}'})
                    doc.close()
                else:
                    tabs_info.append({'index': 0, 'name': '전체'})
        except Exception as e:
            return Response({'found': True, 'error': str(e), 'file': os.path.basename(found_file)})

        return Response({
            'found': True,
            'file': os.path.basename(found_file),
            'ext': ext,
            'school_name': school_name,
            'tabs': tabs_info,
        })

    @action(detail=False, methods=['get'], url_path='nas_file_content')
    def nas_file_content(self, request):
        """NAS 파일 → LibreOffice로 PDF 변환 → 페이지별 이미지 반환"""
        import os, subprocess, base64, hashlib, tempfile, shutil
        from apps.schools.models import School
        school_id = request.query_params.get('school_id')
        tab = request.query_params.get('tab')
        page = int(request.query_params.get('page', 0))
        if not school_id or not tab:
            return Response({'error': 'school_id와 tab 필요'}, status=400)
        try:
            school = School.objects.get(pk=school_id)
        except School.DoesNotExist:
            return Response({'error': '학교 없음'}, status=404)

        folder = self.NAS_FOLDERS.get(tab)
        if not folder:
            return Response({'error': '잘못된 tab'}, status=400)

        nas_dir = os.path.join(self.NAS_BASE, folder)
        found_file = None
        if os.path.isdir(nas_dir):
            for fname in os.listdir(nas_dir):
                if school.name in fname:
                    found_file = os.path.join(nas_dir, fname)
                    break
        if not found_file:
            return Response({'error': '파일 없음'}, status=404)

        ext = found_file.rsplit('.', 1)[-1].lower()

        try:
            # 캐시 디렉토리: 파일 해시 기반
            file_hash = hashlib.md5(found_file.encode() + str(os.path.getmtime(found_file)).encode()).hexdigest()[:12]
            cache_dir = f'/tmp/nms_cache/{file_hash}'

            if ext == 'pdf':
                # PDF → PyMuPDF로 직접 이미지 변환
                import fitz
                doc = fitz.open(found_file)
                if page >= len(doc):
                    doc.close()
                    return Response({'error': '페이지 범위 초과'}, status=400)
                p = doc[page]
                pix = p.get_pixmap(dpi=150)
                img_bytes = pix.tobytes('png')
                doc.close()
                b64 = base64.b64encode(img_bytes).decode()
                return Response({'type': 'image', 'page': page, 'image': f'data:image/png;base64,{b64}'})

            # PPTX/Excel → LibreOffice로 PDF 변환 → PDF → 이미지
            # 캐시된 이미지 확인
            cached_img = os.path.join(cache_dir, f'page_{page}.png')
            if os.path.exists(cached_img):
                with open(cached_img, 'rb') as f:
                    b64 = base64.b64encode(f.read()).decode()
                return Response({'type': 'image', 'page': page, 'image': f'data:image/png;base64,{b64}'})

            # LibreOffice로 PDF 변환
            os.makedirs(cache_dir, exist_ok=True)
            pdf_path = os.path.join(cache_dir, 'converted.pdf')
            if not os.path.exists(pdf_path):
                # LibreOffice 변환
                with tempfile.TemporaryDirectory() as tmpdir:
                    cmd = [
                        'libreoffice', '--headless', '--convert-to', 'pdf',
                        '--outdir', tmpdir, found_file
                    ]
                    subprocess.run(cmd, timeout=120, capture_output=True)
                    # 변환된 PDF 찾기
                    for f in os.listdir(tmpdir):
                        if f.endswith('.pdf'):
                            shutil.copy2(os.path.join(tmpdir, f), pdf_path)
                            break

            if not os.path.exists(pdf_path):
                return Response({'error': 'LibreOffice 변환 실패'}, status=500)

            # PDF → 이미지 (PyMuPDF)
            import fitz
            doc = fitz.open(pdf_path)
            if page >= len(doc):
                doc.close()
                return Response({'error': '페이지 범위 초과'}, status=400)
            p = doc[page]
            pix = p.get_pixmap(dpi=200)
            img_bytes = pix.tobytes('png')
            doc.close()

            # 캐시 저장
            with open(cached_img, 'wb') as f:
                f.write(img_bytes)

            b64 = base64.b64encode(img_bytes).decode()
            return Response({'type': 'image', 'page': page, 'image': f'data:image/png;base64,{b64}'})

        except Exception as e:
            return Response({'error': str(e)}, status=500)
        except Exception as e:
            return Response({'error': str(e)}, status=500)

    @action(detail=False, methods=['get'], url_path='nas_file_download')
    def nas_file_download(self, request):
        """NAS 파일 다운로드"""
        import os
        from urllib.parse import quote
        from django.http import FileResponse
        from apps.schools.models import School
        school_id = request.query_params.get('school_id')
        tab = request.query_params.get('tab')
        if not school_id or not tab:
            return Response({'error': 'school_id와 tab 필요'}, status=400)
        try:
            school = School.objects.get(pk=school_id)
        except School.DoesNotExist:
            return Response({'error': '학교 없음'}, status=404)

        folder = self.NAS_FOLDERS.get(tab)
        if not folder:
            return Response({'error': '잘못된 tab'}, status=400)
        nas_dir = os.path.join(self.NAS_BASE, folder)
        found_file = None
        if os.path.isdir(nas_dir):
            for fname in os.listdir(nas_dir):
                if school.name in fname:
                    found_file = os.path.join(nas_dir, fname)
                    break
        if not found_file:
            return Response({'error': '파일 없음'}, status=404)
        fname = os.path.basename(found_file)
        resp = FileResponse(open(found_file, 'rb'))
        resp['Content-Disposition'] = f"attachment; filename*=UTF-8''{quote(fname)}"
        return resp

    @action(detail=False, methods=['get'])
    def device_counts(self, request):
        """학교별 장비 수량 조회 (정기점검 보고서용) — SchoolEquipment 기준"""
        from apps.schools.models import SchoolEquipment
        school_id = request.query_params.get('school_id')
        if not school_id:
            return Response({'switch_count': 0, 'poe_count': 0, 'ap_count': 0})
        cats = list(SchoolEquipment.objects.filter(
            school_id=school_id
        ).values_list('category', flat=True))
        return Response({
            'switch_count': sum(1 for c in cats if '스위치' in c and 'PoE' not in c),
            'poe_count': sum(1 for c in cats if 'PoE' in c),
            'ap_count': sum(1 for c in cats if 'AP' in c or '무선' in c),
        })

    @action(detail=False, methods=['get'])
    def snmp_guide(self, request):
        """SNMP 설정 가이드 Word 문서 다운로드"""
        import io
        from django.http import HttpResponse
        from apps.schools.models import School
        from .models import NetworkDevice

        school_id = request.query_params.get('school_id')
        if not school_id:
            return Response({'error': 'school_id가 필요합니다.'}, status=status.HTTP_400_BAD_REQUEST)

        try:
            school = School.objects.get(id=school_id)
        except School.DoesNotExist:
            return Response({'error': '학교를 찾을 수 없습니다.'}, status=status.HTTP_404_NOT_FOUND)

        try:
            from docx import Document
            from docx.shared import Pt, RGBColor, Cm
            from docx.enum.text import WD_ALIGN_PARAGRAPH
            from docx.oxml.ns import qn
        except ImportError:
            return Response({'error': 'python-docx가 설치되지 않았습니다.'}, status=status.HTTP_500_INTERNAL_SERVER_ERROR)

        devices = NetworkDevice.objects.filter(school=school).order_by('network_type', 'name')
        doc = Document()

        # ── 한글 폰트 전역 설정 (맑은 고딕) ─────────────────────
        KOREAN_FONT = '맑은 고딕'
        style = doc.styles['Normal']
        style.font.name = KOREAN_FONT
        style.font.size = Pt(10)
        rPr = style.element.get_or_add_rPr()
        rFonts = rPr.find(qn('w:rFonts'))
        if rFonts is None:
            from docx.oxml import OxmlElement
            rFonts = OxmlElement('w:rFonts')
            rPr.append(rFonts)
        rFonts.set(qn('w:ascii'), KOREAN_FONT)
        rFonts.set(qn('w:hAnsi'), KOREAN_FONT)
        rFonts.set(qn('w:eastAsia'), KOREAN_FONT)
        rFonts.set(qn('w:cs'), KOREAN_FONT)
        # 제목 스타일도 한글 폰트 적용
        for heading in ('Heading 1', 'Heading 2', 'Heading 3', 'Title'):
            try:
                h_style = doc.styles[heading]
                h_style.font.name = KOREAN_FONT
                h_rPr = h_style.element.get_or_add_rPr()
                h_rFonts = h_rPr.find(qn('w:rFonts'))
                if h_rFonts is None:
                    from docx.oxml import OxmlElement
                    h_rFonts = OxmlElement('w:rFonts')
                    h_rPr.append(h_rFonts)
                h_rFonts.set(qn('w:ascii'), KOREAN_FONT)
                h_rFonts.set(qn('w:hAnsi'), KOREAN_FONT)
                h_rFonts.set(qn('w:eastAsia'), KOREAN_FONT)
                h_rFonts.set(qn('w:cs'), KOREAN_FONT)
            except KeyError:
                pass

        def _apply_korean_font(run):
            """개별 run에 한글 폰트 명시"""
            run.font.name = KOREAN_FONT
            rPr = run._element.get_or_add_rPr()
            rFonts = rPr.find(qn('w:rFonts'))
            if rFonts is None:
                from docx.oxml import OxmlElement
                rFonts = OxmlElement('w:rFonts')
                rPr.append(rFonts)
            rFonts.set(qn('w:ascii'), KOREAN_FONT)
            rFonts.set(qn('w:hAnsi'), KOREAN_FONT)
            rFonts.set(qn('w:eastAsia'), KOREAN_FONT)
            rFonts.set(qn('w:cs'), KOREAN_FONT)

        # 제목
        title = doc.add_heading(f'{school.name} SNMP 설정 가이드', 0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        for run in title.runs:
            _apply_korean_font(run)

        date_p = doc.add_paragraph(f'작성일: {timezone.now().strftime("%Y년 %m월 %d일")}')
        for run in date_p.runs:
            _apply_korean_font(run)
        doc.add_paragraph()

        def _add_para(text, style=None, bold=False):
            """한글 폰트 자동 적용 paragraph 생성"""
            if style:
                p = doc.add_paragraph(style=style)
                run = p.add_run(text)
            else:
                p = doc.add_paragraph()
                run = p.add_run(text)
            _apply_korean_font(run)
            if bold:
                run.bold = True
            return p

        def _add_heading(text, level):
            h = doc.add_heading(text, level=level)
            for run in h.runs:
                _apply_korean_font(run)
            return h

        # 1. 개요
        _add_heading('1. SNMP 개요', 1)
        _add_para(
            'SNMP(Simple Network Management Protocol)는 네트워크 장비의 상태를 모니터링하기 위한 표준 프로토콜입니다. '
            '본 시스템에서는 SNMPv2c를 사용하여 장비의 가동 상태, 트래픽, 포트 상태를 수집합니다.'
        )

        # 2. 설정 방법
        _add_heading('2. 장비별 SNMP 설정 방법', 1)
        _add_heading('2-1. CBS/C3100/C3500 시리즈 (코어/분배 스위치)', 2)
        for cmd in [
            'snmp-server community public RO',
            'snmp-server community private RW',
            'snmp-server enable traps',
            'snmp-server host [NMS서버IP] traps public',
        ]:
            p = doc.add_paragraph(style='List Bullet')
            run = p.add_run(cmd)
            run.font.name = 'Courier New'

        _add_heading('2-2. GS724T / SG300 시리즈 (접속 스위치)', 2)
        _add_para('웹 관리 인터페이스 접속 → Security → SNMP → Communities 메뉴에서 설정')
        for step in ['Community String: public (Read Only)', 'Trap Host: [NMS서버IP]', 'SNMP Version: v2c']:
            _add_para(step, style='List Bullet')

        # 3. 장비 목록
        _add_heading('3. 장비 목록 및 설정 현황', 1)

        table = doc.add_table(rows=1, cols=6)
        table.style = 'Table Grid'
        hdr = table.rows[0].cells
        for i, h in enumerate(['장비명', '모델', '설치위치', '망구분', 'IP주소', 'SNMP']):
            hdr[i].text = h
            for run in hdr[i].paragraphs[0].runs:
                _apply_korean_font(run)
                run.font.bold = True

        TYPE_KO = {'switch':'스위치','poe_switch':'PoE스위치','ap':'AP','router':'라우터','firewall':'방화벽','server':'서버'}
        for d in devices:
            row = table.add_row().cells
            values = [
                d.name,
                d.model or '-',
                d.location or '-',
                d.network_type or '-',
                d.ip_address or '미등록',
                '설정완료' if d.snmp_enabled else '미설정',
            ]
            for i, v in enumerate(values):
                row[i].text = v
                for run in row[i].paragraphs[0].runs:
                    _apply_korean_font(run)

        # 4. NMS 연동 절차
        _add_heading('4. NMS 연동 절차', 1)
        steps = [
            '장비에 SNMP Community String 설정 (public/private)',
            '장비 IP 주소를 NMS 시스템에 등록',
            '자산 관리 → 네트워크 설정 → SNMP 활성화 체크',
            'NMS 모니터링 탭에서 장비 상태 확인',
            '장애 발생 시 이벤트 알림 자동 수신',
        ]
        for i, s in enumerate(steps, 1):
            _add_para(f'{i}. {s}')

        buf = io.BytesIO()
        doc.save(buf)
        buf.seek(0)

        from datetime import datetime
        import urllib.parse
        today = datetime.now().strftime('%Y%m%d')
        response = HttpResponse(
            buf.read(),
            content_type='application/vnd.openxmlformats-officedocument.wordprocessingml.document'
        )
        filename = f'SNMP설정가이드_{school.name}_{today}.docx'
        encoded = urllib.parse.quote(filename)
        response['Content-Disposition'] = f"attachment; filename*=UTF-8''{encoded}"

        # NAS 자동 저장
        try:
            from .services import write_topology_files_to_nas
            write_topology_files_to_nas(school)
        except Exception:
            pass
        return response

    @action(detail=False, methods=['post'])
    def scan_pptx(self, request):
        """PPTX 구성도 스캔 (NAS /산출물/{school_id}/구성도/*.pptx)

        body:
          {"school_id": 123}   개별 학교만 동기 실행
          {}                    전체 학교 비동기 실행 (Celery)
        """
        from .tasks import scan_network_pptx
        school_id = request.data.get('school_id')
        try:
            if school_id:
                # 개별 학교는 즉시 실행하고 결과 반환
                result = scan_network_pptx(school_id=int(school_id))
                return Response(result, status=status.HTTP_200_OK)
            # 전체는 비동기
            task = scan_network_pptx.delay(school_id=None)
            return Response({'status': 'started', 'task_id': str(task.id)},
                            status=status.HTTP_202_ACCEPTED)
        except Exception as e:
            return Response({'error': str(e)}, status=status.HTTP_500_INTERNAL_SERVER_ERROR)

    @action(detail=False, methods=['get'])
    def pptx_status(self, request):
        """특정 학교의 PPTX 파일/파싱 상태 조회"""
        import os
        school_id = request.query_params.get('school_id')
        if not school_id:
            return Response({'error': 'school_id 필요'}, status=status.HTTP_400_BAD_REQUEST)

        base_dir = os.environ.get('NAS_ARTIFACT_ROOT', '/app/nas/media/npms/산출물')
        pptx_dir = os.path.join(base_dir, str(school_id), '구성도')
        files = []
        if os.path.isdir(pptx_dir):
            for f in os.listdir(pptx_dir):
                if f.lower().endswith('.pptx') and not f.startswith('.'):
                    fp = os.path.join(pptx_dir, f)
                    files.append({
                        'filename': f,
                        'size': os.path.getsize(fp),
                        'mtime': os.path.getmtime(fp),
                    })

        topo = NetworkTopology.objects.filter(school_id=school_id).first()
        return Response({
            'school_id': int(school_id),
            'pptx_dir': pptx_dir,
            'pptx_files': sorted(files, key=lambda x: -x['mtime']),
            'topology_exists': topo is not None,
            'slide_titles': topo.slide_titles if topo else [],
            'pptx_path': topo.pptx_path if topo else '',
            'pptx_mtime': topo.pptx_mtime.isoformat() if topo and topo.pptx_mtime else None,
            'updated_at': topo.updated_at.isoformat() if topo else None,
        })


class NetworkEventViewSet(viewsets.ModelViewSet):
    """네트워크 이벤트/알림"""
    serializer_class = NetworkEventSerializer
    permission_classes = [IsAuthenticated]
    pagination_class = StandardPagination

    def get_queryset(self):
        qs = NetworkEvent.objects.select_related('device', 'device__school')
        resolved = self.request.query_params.get('resolved')
        if resolved == '0':
            qs = qs.filter(is_resolved=False)
        elif resolved == '1':
            qs = qs.filter(is_resolved=True)
        severity = self.request.query_params.get('severity')
        if severity:
            qs = qs.filter(severity=severity)
        school_id = self.request.query_params.get('school_id')
        if school_id:
            qs = qs.filter(device__school_id=school_id)
        return qs

    @action(detail=True, methods=['post'])
    def resolve(self, request, pk=None):
        """이벤트 해결 처리"""
        event = self.get_object()
        if event.is_resolved:
            return Response({'error': '이미 해결된 이벤트입니다.'}, status=status.HTTP_400_BAD_REQUEST)
        event.is_resolved = True
        event.resolved_at = timezone.now()
        event.save(update_fields=['is_resolved', 'resolved_at'])
        return Response(NetworkEventSerializer(event).data)

    @action(detail=False, methods=['get'])
    def active_summary(self, request):
        """미해결 이벤트 요약"""
        qs = NetworkEvent.objects.filter(is_resolved=False)
        by_severity = qs.values('severity').annotate(cnt=Count('id'))
        return Response({
            'total': qs.count(),
            'by_severity': {item['severity']: item['cnt'] for item in by_severity},
            'latest': NetworkEventSerializer(qs.order_by('-occurred_at')[:5], many=True).data,
        })


class NetworkCommandViewSet(viewsets.ReadOnlyModelViewSet):
    """원격 명령 실행 이력"""
    serializer_class = NetworkCommandSerializer
    permission_classes = [IsAuthenticated]
    pagination_class = StandardPagination

    def get_queryset(self):
        qs = NetworkCommand.objects.select_related('device', 'executed_by')
        device_id = self.request.query_params.get('device_id')
        if device_id:
            qs = qs.filter(device_id=device_id)
        st = self.request.query_params.get('status')
        if st:
            qs = qs.filter(status=st)
        return qs
